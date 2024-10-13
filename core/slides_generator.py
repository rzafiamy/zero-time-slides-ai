from pptx import Presentation
from core.image_search import ImageSearcher
from core.ai_requester import AIRequester
from core.layout_manager import LayoutManager
from core.web_search import WebSearch
import os

class SlidesGenerator:
    def __init__(self, config):
        """
        Initialize the SlidesGenerator with the config for paths and image search settings.
        """
        self.config = config
        self.ai_requester = AIRequester(config)
        self.layoutM = LayoutManager()

    def generate_presentation(self, topic, slide_length):
        """
        Generate a PowerPoint presentation based on the AI content response.
        """
        context = WebSearch(self.config, topic).perform_search_and_format(self.ai_requester)
        prompt = self.ai_requester.create_prompt(topic, slide_length, context)
        content = self.ai_requester.request_ai(prompt)
        
        presentation = Presentation("templates/template0.pptx")
        self._parse_response(presentation, content)
        title = self._get_presentation_title(presentation) or topic
        ppt_filename = f"{title}.pptx"
        presentation.save(f"{self.config['OUTPUT_PATH']}/{ppt_filename}")
        return ppt_filename

    def _parse_response(self, presentation, content):
        """
        Parse the response from the AI model and create slides in the PowerPoint presentation.
        """
        slides = content.split("[SLIDEBREAK]")
        for slide in slides:

            slide_type = self._find_slide_type(slide)
            title = self._extract_tag_content(slide, "[TITLE]", "[/TITLE]") or "Untitled"
            

            if slide_type == "[L_TS]":  # Title Slide
                subtitle = self._extract_tag_content(slide, "[SUBTITLE]", "[/SUBTITLE]") 
                self.layoutM._create_title_slide(presentation, title, subtitle)
            elif slide_type == "[L_CS]":  # Content Slide
                context = self._extract_tag_content(slide, "[CONTENT]", "[/CONTENT]")
                prompt =  self.ai_requester.create_slide_prompt(slide_type, title, context)
                text_content =  self.ai_requester.request_ai(prompt)

                self.layoutM._create_title_and_content_slide(presentation, title, text_content)
            elif slide_type == "[L_IS]":  # Image Slide
                context = self._extract_tag_content(slide, "[CONTENT]", "[/CONTENT]")
                prompt = self.ai_requester.create_slide_prompt(slide_type, title, context)
                text_content =  self.ai_requester.request_ai(prompt)


                image_query = self._extract_tag_content(slide, "[IMAGE]", "[/IMAGE]")
                image_searcher = ImageSearcher(self.config)
                image_path = image_searcher.download_image(title)
        
                self.layoutM._create_picture_with_caption_slide(presentation, title, 
                    os.path.join(self.config["IMAGES_PATH"], image_path),
                    text_content)
            elif slide_type == "[L_TCS]":  # Two Content Slide
                contents = self._extract_all_tag_content(slide, "[CONTENT]", "[/CONTENT]")
                if len(contents) < 2:
                    continue
                content_left = contents[0]
                content_right = contents[1]

                prompt_left =  self.ai_requester.create_slide_prompt(slide_type, title, content_left)
                text_content_left =  self.ai_requester.request_ai(prompt_left)

                prompt_right =  self.ai_requester.create_slide_prompt(slide_type, title, content_right)
                text_content_right =  self.ai_requester.request_ai(prompt_right)

                self.layoutM._create_two_content_slide(presentation, title, text_content_left, text_content_right)
            elif slide_type == "[L_THS]":  # Thanks Slide
                self.layoutM._create_title_only_slide(presentation, title)

    
    def _extract_all_tag_content(self, text, start_tag, end_tag):
        """
        Extract all occurrences of content between the custom start and end tags in the given text.
        """
        contents = []
        start = 0

        while True:
            start = text.find(start_tag, start)
            if start == -1:  # No more tags
                break
            end = text.find(end_tag, start + len(start_tag))
            if end == -1:  # End tag not found
                break

            # Extract content and add to list
            content = text[start + len(start_tag):end].strip()
            contents.append(content)

            # Move start position beyond the current tag
            start = end + len(end_tag)

        return contents


    def _extract_tag_content(self, text, start_tag, end_tag):
        """
        Extract the content between the start and end tags in the given text.
        """
        start = text.find(start_tag)
        end = text.find(end_tag)
        if start != -1 and end != -1:
            return text[start + len(start_tag):end].strip()
        return ""

    def _find_slide_type(self, text):
        """
        Find the slide type tag in the given text.
        """
        for tag in ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]", "[L_TCS]"]:
            if tag in text:
                return tag
        return None

    def _get_presentation_title(self, presentation):
        """
        Get the title of the presentation from the first slide.
        """
        return presentation.slides[0].shapes.title.text
