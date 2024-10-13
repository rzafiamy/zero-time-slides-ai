import argparse
from pptx import Presentation

# Function to list all layouts and placeholders
def list_layouts_and_placeholders(pptx_file):
    # Open the presentation
    presentation = Presentation(pptx_file)
    
    print(f"Total number of slide layouts: {len(presentation.slide_layouts)}\n")
    
    # Iterate through each slide layout in the presentation
    for idx, layout in enumerate(presentation.slide_layouts):
        print(f"Slide Layout {idx}: {layout.name}")
        
        # Get all placeholders in the layout
        placeholders = layout.placeholders
        if placeholders:
            for placeholder in placeholders:
                print(f"\tPlaceholder {placeholder.placeholder_format.idx}: {placeholder.name}, Type: {placeholder.placeholder_format.type}")
        else:
            print("\tNo placeholders in this layout")
        print("\n")

# Main function to handle argparse and calling the function
def main():
    # Set up argument parser
    parser = argparse.ArgumentParser(description='List all slide layouts and placeholders in a PPTX file.')
    parser.add_argument('pptx_file', type=str, help='Path to the PPTX file')

    # Parse the command-line arguments
    args = parser.parse_args()

    # Call the function with the provided pptx file
    list_layouts_and_placeholders(args.pptx_file)

# Execute the script
if __name__ == "__main__":
    main()
